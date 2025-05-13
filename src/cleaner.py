import tkinter as tk
from tkinter import filedialog, colorchooser, messagebox
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import io
import pytesseract
from PIL import Image
import shutil
import re

# Search for the tesseract executable in the system PATH
tesseract_path = shutil.which("tesseract")

if tesseract_path is None:
    raise EnvironmentError("Tesseract executable not found in PATH. Please install Tesseract OCR.")
else:
    pytesseract.pytesseract.tesseract_cmd = tesseract_path

def hex_to_rgb_color(hex_color):
    """Convert a hex color (e.g. '#FF0000') to an RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def is_mostly_text(s):
    """Heuristic: if the string is long enough and the ratio of letters is high, treat it as text."""
    s = s.strip()
    if len(s) < 5:
        return False
    alpha_count = sum(c.isalpha() for c in s)
    return (alpha_count / len(s)) > 0.5

# Define text spacing options with their corresponding numeric values
# Values are in points/100 (e.g. -1.5 points = -150)
TEXT_SPACING_OPTIONS = {
    "Very Tight": -150,  # -1.5 points
    "Tight": -50,        # -0.5 points
    "Normal": 0,         # 0 points (default)
    "Loose": 50,         # 0.5 points
    "Very Loose": 150    # 1.5 points
}

class PPTCleanerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Cleaner")
        self.input_file = None
        self.output_file = None
        self.text_color_hex = "#000000"  # default text color: black
        self.bg_color_hex = "#FFFFFF"    # default background color: white
        self.setup_gui()

    def setup_gui(self):
        # File selection frame
        file_frame = tk.Frame(self.root, padx=10, pady=10)
        file_frame.pack(fill="x")
        tk.Label(file_frame, text="Select PPTX File:").grid(row=0, column=0, sticky="w")
        self.input_entry = tk.Entry(file_frame, width=50)
        self.input_entry.grid(row=0, column=1, padx=5)
        tk.Button(file_frame, text="Browse", command=self.browse_file).grid(row=0, column=2)

        # Output file selection frame
        out_frame = tk.Frame(self.root, padx=10, pady=10)
        out_frame.pack(fill="x")
        tk.Label(out_frame, text="Save Processed File as:").grid(row=0, column=0, sticky="w")
        self.output_entry = tk.Entry(out_frame, width=50)
        self.output_entry.grid(row=0, column=1, padx=5)
        tk.Button(out_frame, text="Browse", command=self.browse_output).grid(row=0, column=2)

        # Settings frame
        settings_frame = tk.LabelFrame(self.root, text="Settings", padx=10, pady=10)
        settings_frame.pack(fill="x", padx=10, pady=5)

        # Custom Font Enable/Disable
        self.custom_font_enabled_var = tk.IntVar(value=0)
        self.custom_font_checkbox = tk.Checkbutton(settings_frame, text="Enable Custom Font", variable=self.custom_font_enabled_var, command=self.toggle_font_options)
        self.custom_font_checkbox.grid(row=0, column=0, sticky="w", pady=2)

        # Font family dropdown
        tk.Label(settings_frame, text="Font Family:").grid(row=1, column=0, sticky="w")
        font_families = ["Calibri", "Aptos", "Arial", "Times New Roman", "Verdana", "Helvetica"]
        self.font_family_var = tk.StringVar(value="Calibri")
        self.font_family_dropdown = tk.OptionMenu(settings_frame, self.font_family_var, *font_families)
        self.font_family_dropdown.config(state="disabled")
        self.font_family_dropdown.grid(row=1, column=1, padx=5, pady=2, sticky="w")

        # Font size dropdown
        tk.Label(settings_frame, text="Font Size (pt):").grid(row=2, column=0, sticky="w")
        font_sizes = ["12", "14", "16", "18", "20", "22", "24", "26", "28", "30"]
        self.font_size_var = tk.StringVar(value="24")
        self.font_size_dropdown = tk.OptionMenu(settings_frame, self.font_size_var, *font_sizes)
        self.font_size_dropdown.config(state="disabled")
        self.font_size_dropdown.grid(row=2, column=1, padx=5, pady=2, sticky="w")

        # Text spacing dropdown - NEW FEATURE
        tk.Label(settings_frame, text="Text Spacing:").grid(row=3, column=0, sticky="w")
        spacing_options = list(TEXT_SPACING_OPTIONS.keys())
        self.text_spacing_var = tk.StringVar(value="Normal")
        self.text_spacing_dropdown = tk.OptionMenu(settings_frame, self.text_spacing_var, *spacing_options)
        self.text_spacing_dropdown.grid(row=3, column=1, padx=5, pady=2, sticky="w")

        # Text formatting settings
        self.bold_var = tk.IntVar(value=1)
        tk.Checkbutton(settings_frame, text="Bold Text", variable=self.bold_var).grid(row=4, column=0, sticky="w", pady=2)

        tk.Label(settings_frame, text="Text Color:").grid(row=5, column=0, sticky="w")
        self.text_color_button = tk.Button(settings_frame, bg=self.text_color_hex, width=3, command=self.choose_text_color)
        self.text_color_button.grid(row=5, column=1, sticky="w", padx=5, pady=2)

        # Duplicate text removal setting
        self.dup_var = tk.IntVar(value=1)
        tk.Checkbutton(settings_frame, text="Remove Duplicate Textboxes per Slide", variable=self.dup_var)\
            .grid(row=6, column=0, columnspan=2, sticky="w", pady=2)

        # Background color setting
        tk.Label(settings_frame, text="Slide Background Color:").grid(row=7, column=0, sticky="w")
        self.bg_color_button = tk.Button(settings_frame, bg=self.bg_color_hex, width=3, command=self.choose_bg_color)
        self.bg_color_button.grid(row=7, column=1, sticky="w", padx=5, pady=2)

        # Additional cleanup option (placeholder for removing animations)
        self.remove_animations_var = tk.IntVar(value=0)
        tk.Checkbutton(settings_frame, text="Remove Animations (if possible)", variable=self.remove_animations_var)\
            .grid(row=8, column=0, columnspan=2, sticky="w", pady=2)

        # Enable OCR on images
        self.enable_ocr_var = tk.IntVar(value=0)
        tk.Checkbutton(settings_frame, text="Enable OCR on Images", variable=self.enable_ocr_var)\
            .grid(row=9, column=0, columnspan=2, sticky="w", pady=2)
            
        # Remove presentation theme
        self.remove_theme_var = tk.IntVar(value=0)
        tk.Checkbutton(settings_frame, text="Remove Presentation Theme", variable=self.remove_theme_var)\
            .grid(row=10, column=0, columnspan=2, sticky="w", pady=2)

        # Process button
        process_frame = tk.Frame(self.root, padx=10, pady=10)
        process_frame.pack(fill="x")
        tk.Button(process_frame, text="Process PPTX", command=self.process_file).pack()

    def toggle_font_options(self):
        state = "normal" if self.custom_font_enabled_var.get() else "disabled"
        self.font_family_dropdown.config(state=state)
        self.font_size_dropdown.config(state=state)

    def browse_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.input_file = file_path
            self.input_entry.delete(0, tk.END)
            self.input_entry.insert(0, file_path)

    def browse_output(self):
        file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.output_file = file_path
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, file_path)

    def choose_text_color(self):
        color_code = colorchooser.askcolor(title="Choose text color")
        if color_code[1]:
            self.text_color_hex = color_code[1]
            self.text_color_button.config(bg=self.text_color_hex)

    def choose_bg_color(self):
        color_code = colorchooser.askcolor(title="Choose background color")
        if color_code[1]:
            self.bg_color_hex = color_code[1]
            self.bg_color_button.config(bg=self.bg_color_hex)

    def process_file(self):
        if not self.input_file or not os.path.isfile(self.input_file):
            messagebox.showerror("Error", "Please select a valid PPTX file.")
            return
        if not self.output_file:
            messagebox.showerror("Error", "Please select an output file location.")
            return

        try:
            # Determine if custom font settings are enabled
            custom_font_enabled = bool(self.custom_font_enabled_var.get())
            if custom_font_enabled:
                custom_font = self.font_family_var.get()
                custom_font_size = Pt(int(self.font_size_var.get()))
            else:
                custom_font = None
                custom_font_size = None

            # Get the selected text spacing option
            text_spacing_value = TEXT_SPACING_OPTIONS[self.text_spacing_var.get()]

            settings = {
                'enable_custom_font': custom_font_enabled,
                'custom_font': custom_font,
                'custom_font_size': custom_font_size,
                'text_spacing': text_spacing_value,  # Add text spacing setting
                'text_bold': bool(self.bold_var.get()),
                'text_color': hex_to_rgb_color(self.text_color_hex),
                'remove_duplicates': bool(self.dup_var.get()),
                'background_color': hex_to_rgb_color(self.bg_color_hex),
                'remove_animations': bool(self.remove_animations_var.get()),
                'enable_ocr': bool(self.enable_ocr_var.get()),
                'remove_wordart': False,  # Default setting
                'remove_theme': bool(self.remove_theme_var.get())  # New theme removal setting
            }

            self.process_pptx(self.input_file, self.output_file, settings)
            messagebox.showinfo("Success", f"Processed file saved as:\n{self.output_file}")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")

    def process_pptx(self, input_path, output_path, settings):
        prs = Presentation(input_path)

        # Define simple word replacements
        replacements = {
            'bvb': 'bv',
            'pt': 'PT',
            'pten': 'PTen',
            'pte': 'PTe',
            'wrs': 'wss',
            'knn': 'kunnen',
            'vr': 'voor',
            'wilt': 'wil',
            'versch': 'verschillend',
        }

        # Regex pattern will be built dynamically per word
        for slide in prs.slides:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        for abbr, full in replacements.items():
                            # Match: space before, space/dot/comma/semicolon after (lookahead), case-insensitive
                            pattern = rf'(?<=\s){abbr}(?=[\s\.,;])'
                            text = re.sub(pattern, full, text, flags=re.IGNORECASE)
                        shape.text_frame.text = text

        # unicode_replacements = {
        #     '\uf075': '\u2022',  # Replace unicode character with bullet
        #     '\uf075\u0009': '\u2022  ',  # Replace unicode character with bullet
        #     '\uf075\u0009\u0009': '\u2022    ',  # Replace unicode character with bullet
        #     '\uf0e7': '\u2022',  # Another unicode character to replace
        # }

        # # Edit all text in the presentation to replace unicode characters
        # for slide in prs.slides:
        #     for shape in slide.shapes:
        #         if shape.has_text_frame:
        #             for paragraph in shape.text_frame.paragraphs:
        #                 for run in paragraph.runs:
        #                     for unicode_char, replacement in unicode_replacements.items():
        #                         if unicode_char in run.text:
        #                             run.text = run.text.replace(unicode_char, replacement)

        # for slide in prs.slides:
        #     for shape in slide.shapes:
        #         if shape.has_text_frame:
        #             for paragraph in shape.text_frame.paragraphs:
        #                 if paragraph.text:
        #                     print(paragraph.text)  # Debug before

        #                     # Use regex to find first alphabetic character
        #                     match = re.search(r'([^\w]*)(\w)(.*)', paragraph.text, re.UNICODE)
        #                     if match:
        #                         prefix, first_letter, rest = match.groups()
        #                         paragraph.text = f"{prefix}{first_letter.upper()}{rest}"

        #                     print(paragraph.text)  # Debug after


        
        # # Remove presentation theme if requested
        # if settings['remove_theme']:
        #     # Remove theme by clearing all theme parts
        #     # This removes all theme-based formatting that might conflict with our custom settings
        #     for i in range(len(prs.part.part_related_by_rel_type('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'))):
        #         try:
        #             # Access theme part by index since the direct reference may change during deletion
        #             theme_part = prs.part.part_related_by_rel_type('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
        #             if theme_part:
        #                 # Get the relationship ID
        #                 rel_ids = [rel_id for rel_id, part in prs.part.rels.items() 
        #                            if part.target_part == theme_part]
        #                 if rel_ids:
        #                     # Remove the relationship to the theme part
        #                     prs.part.drop_rel(rel_ids[0])
        #         except Exception as e:
        #             print(f"Error removing theme part: {e}")

        # # Iterate through each slide master
        # for master in prs.slide_masters:
        #     for shape in master.shapes:
        #         # Check if the shape is a picture or an autoshape
        #         if shape.shape_type in [1, 13]:  # 1: Picture, 13: AutoShape
        #             # Hide the shape
        #             shape.element.getparent().remove(shape.element)
        #                 # Process each slide
        # for slide in prs.slides:
        #     # Remove theme reference at slide level if requested
        #     if settings['remove_theme']:
        #         try:
        #             # Remove slide layout relationship which contains theme-based formatting
        #             if hasattr(slide, '_element') and hasattr(slide._element, 'rId'):
        #                 slide_element = slide._element
        #                 for rel in list(slide_element.rels.values()):
        #                     if 'slideLayout' in rel.reltype:
        #                         slide_element.drop_rel(rel.rId)
        #         except Exception as e:
        #             print(f"Error removing slide-level theme reference: {e}")
                    
        #     seen_texts = set()
        #     shapes_to_remove = []
        #     # Process shapes for text cleanup and OCR conversion
        #     for shape in slide.shapes:
        #         # Process text shapes (if they have text)
        #         if shape.has_text_frame:
        #             text = shape.text.strip()
        #             # Remove duplicate textboxes if enabled
        #             if settings['remove_duplicates']:
        #                 if text in seen_texts:
        #                     shapes_to_remove.append(shape)
        #                     continue
        #                 seen_texts.add(text)

        #             if ("WordArt" in shape.name and settings['remove_wordart']) or "WordArt" not in shape.name:
        #                 # Apply text formatting for regular text shapes
        #                 for paragraph in shape.text_frame.paragraphs:
        #                     for run in paragraph.runs:
        #                         if settings['enable_custom_font'] and settings['custom_font']:
        #                             run.font.name = settings['custom_font']
        #                             # if run.font.size:
        #                                 # run.font.size = min(settings['custom_font_size'], run.font.size)
        #                             # else:
        #                             run.font.size = settings['custom_font_size']
        #                         run.font.bold = settings['text_bold']
        #                         run.font.color.rgb = settings['text_color']
        #                         # Apply text spacing setting - NEW FEATURE
        #                         # Character spacing is set in space-per-100 units
        #                         if hasattr(run.font, 'char_spacing'):
        #                             run.font.char_spacing = settings['text_spacing']
        #                         if run.font.underline:
        #                             run.font.color.rgb = settings['text_color']

        #         # Process picture shapes for OCR if enabled
        #         elif settings['enable_ocr'] and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #             try:
        #                 image_stream = io.BytesIO(shape.image.blob)
        #                 image = Image.open(image_stream)
        #                 ocr_text = pytesseract.image_to_string(image).strip()
        #                 if is_mostly_text(ocr_text):
        #                     # Add a new textbox with OCR text at the same position and size
        #                     left = shape.left
        #                     top = shape.top
        #                     width = shape.width
        #                     height = shape.height
        #                     textbox = slide.shapes.add_textbox(left, top, width, height)
        #                     tf = textbox.text_frame
        #                     tf.clear()  # remove any default paragraph
        #                     p = tf.add_paragraph()
        #                     p.text = ocr_text
        #                     run = p.runs[0] if p.runs else p.add_run()
        #                     if settings['enable_custom_font'] and settings['custom_font']:
        #                         run.font.name = settings['custom_font']
        #                         run.font.size = settings['custom_font_size']
        #                     run.font.bold = settings['text_bold']
        #                     run.font.color.rgb = settings['text_color']
        #                     # Apply text spacing to OCR text - NEW FEATURE
        #                     if hasattr(run.font, 'char_spacing'):
        #                         run.font.char_spacing = settings['text_spacing']
        #                     # Mark the image shape for removal
        #                     shapes_to_remove.append(shape)
        #             except Exception as e:
        #                 print("Error processing image for OCR:", e)

        #     # Remove marked duplicate textboxes and images replaced by OCR text
        #     for shape in shapes_to_remove:
        #         try:
        #             sp = shape._element
        #             sp.getparent().remove(sp)
        #         except Exception as e:
        #             print("Error removing shape:", e)

        #     # Change the slide background color
        #     try:
        #         # Access the slide's background fill
        #         background = slide.background
        #         fill = background.fill
        #         fill.solid()
        #         fill.fore_color.rgb = settings['background_color']
        #     except Exception as e:
        #         print("Error setting background:", e)

        #     # Optional: Remove animations (if supported)
        #     if settings['remove_animations']:
        #         # python-pptx has limited support for animations.
        #         # This is a placeholder for additional cleanup.
        #         pass

        prs.save(output_path)

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCleanerApp(root)
    root.mainloop()