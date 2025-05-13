#!/usr/bin/env python3
import os
import argparse

from pptx import Presentation
from pptx.util import Pt

def clean_spacing(pptx_path):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                # Clear paragraph spacing
                # paragraph.space_before = Pt(0)
                # paragraph.space_after  = Pt(0)
                # paragraph.line_spacing = Pt(0)
                # Clear run-level character spacing (kerning)
                for run in paragraph.runs:
                    run.font.kerning = Pt(0)
    # Save with _clean suffix
    base, ext = os.path.splitext(pptx_path)
    out_path = f"{base}_clean{ext}"
    prs.save(out_path)
    print(f"Saved cleaned presentation as: {out_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Remove all paragraph and run spacing in a PPTX"
    )
    parser.add_argument(
        "pptx_file", help="Path to the source .pptx file"
    )
    args = parser.parse_args()
    clean_spacing(args.pptx_file)
